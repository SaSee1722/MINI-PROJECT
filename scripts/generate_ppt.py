from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def add_bulleted_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = title_text
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        p.font.size = Pt(18)
        p.level = 0
    return slide

def add_link_slide(prs, title_text, intro_lines, link_text, link_url):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = title_text
    tf = body.text_frame
    tf.clear()
    # Intro lines
    for i, line in enumerate(intro_lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        p.font.size = Pt(18)
        p.level = 0
    # Link paragraph
    lp = tf.add_paragraph()
    lp.font.size = Pt(18)
    lp.level = 0
    run = lp.add_run()
    run.text = link_text
    try:
        run.hyperlink.address = link_url
    except Exception:
        pass
    lp.alignment = PP_ALIGN.LEFT
    return slide

def build_presentation(output_path):
    prs = Presentation()

    # 1. Title Page
    add_bulleted_slide(prs, "Attendance App", [
        "Role-based, period-wise attendance with reports",
        "Tech Stack: React + Vite, TailwindCSS, Supabase (Auth, Postgres), jsPDF",
        "Roles: Admin, Staff",
        "Deployment: SPA-ready (Netlify/Vercel)"
    ])

    # 2. Abstract
    add_bulleted_slide(prs, "Abstract", [
        "Problem: manual attendance is slow, error-prone, and fragmented",
        "Solution: cloud-backed web app with secure auth and RBAC",
        "Features: period-wise student + staff tracking, PDF exports",
        "Outcome: faster marking, centralized records, audit-ready reports"
    ])

    # 3. Introduction
    add_bulleted_slide(prs, "Introduction", [
        "Motivation: reduce errors and improve visibility across departments",
        "Objectives: auth + RBAC, interactive timetable, reporting",
        "Scope: React SPA with Supabase; focuses on period tracking"
    ])

    # 4. Literature Survey
    add_bulleted_slide(prs, "Literature Survey", [
        "Manual registers: low-cost; poor aggregation and analytics",
        "Biometric/RFID: accurate; higher cost and maintenance",
        "QR-based: fast; device coordination challenges",
        "Cloud web systems: centralized, scalable; needs solid auth/data design",
        "Gap solved: RBAC, period granularity, bulk ops, integrated reports"
    ])

    # 5. Proposed Work
    add_bulleted_slide(prs, "Proposed Work", [
        "Architecture: React SPA → Supabase → Postgres tables",
        "Auth & RBAC: session + ProtectedRoute with role checks",
        "Period-wise attendance: mark, counts, alt staff, per-student statuses",
        "Bulk imports and PDF reporting"
    ])

    # 6. Flowcharts & Algorithm
    add_bulleted_slide(prs, "Flowcharts & Algorithm", [
        "Auth flow: sign up/in → session → role-based routing",
        "Marking flow: timetable → select period → upsert records",
        "Data write: period_attendance + period_student_attendance",
        "Export flow: generate Period-wise PDF"
    ])

    # 7. Future Work
    add_bulleted_slide(prs, "Future Work", [
        "Offline-first PWA and sync queue",
        "Mobile app (Android/iOS)",
        "Face recognition integration",
        "Notifications/approvals and analytics dashboards"
    ])

    # 8. Results & Conclusions
    add_bulleted_slide(prs, "Results & Conclusions", [
        "Achieved: secure RBAC, interactive timetable, period-wise marking",
        "Reports: color-coded PDFs for audits",
        "Impact: faster cycles and centralized data",
        "Conclusion: practical, scalable attendance management"
    ])

    # 9. Project Demo with hyperlink
    add_link_slide(
        prs,
        "Project Demo",
        [
            "Admin: mark period-wise attendance and download reports",
            "Staff: mark own attendance and class-wise summaries"
        ],
        "Open Demo: mini-project-salabadee.vercel.app",
        "https://mini-project-salabadee.vercel.app/"
    )

    prs.save(output_path)

if __name__ == "__main__":
    build_presentation("docs/Attendance_App_Demo.pptx")